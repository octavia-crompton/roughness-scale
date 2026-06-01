#!/usr/bin/env python3
"""Build a PowerPoint deck with registered figures (+ title & setup slides).

Equations are in latex/equations.tex — not in the PowerPoint.
Figures are sized to fill each slide as large as possible.
"""
import pathlib, re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN

FIG_DIR = pathlib.Path(__file__).resolve().parent.parent / "figures" / "runaround_smooth"
CONCISE = FIG_DIR / "figure_registry_concise.txt"
OUT     = FIG_DIR / "roughness_scale_figures.pptx"

# ═══════════════════════════════════════════════════════════════════════════
# Slide helpers
# ═══════════════════════════════════════════════════════════════════════════

def _textbox(slide, left, top, w, h, text="", sz=18, bold=False, italic=False,
             align=PP_ALIGN.LEFT, color=None):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(sz)
    p.font.bold = bold; p.font.italic = italic; p.alignment = align
    if color:
        from pptx.dml.color import RGBColor
        p.font.color.rgb = RGBColor(*color)
    return tf

def _bullet(tf, text, sz=16, bold=False, level=0):
    p = tf.add_paragraph(); p.text = text; p.font.size = Pt(sz)
    p.font.bold = bold; p.level = level
    return p

# ═══════════════════════════════════════════════════════════════════════════
# Parse concise registry
# ═══════════════════════════════════════════════════════════════════════════
entries = []
with open(CONCISE) as f:
    rlines = f.readlines()
i = 0
while i < len(rlines):
    line = rlines[i].strip()
    mt = re.match(r'^(Fig \d+|SI\d+)\s*[—–-]\s*(\S+)', line)
    if mt:
        fig_id, filename = mt.group(1), mt.group(2)
        cap = []; i += 1
        while i < len(rlines) and (rlines[i].startswith('  ') or rlines[i].strip() == ''):
            s = rlines[i].strip()
            if s: cap.append(s)
            i += 1
        entries.append((fig_id, filename, ' '.join(cap)))
    else:
        i += 1

main_figs = [e for e in entries if e[0].startswith('Fig')]
si_figs   = [e for e in entries if e[0].startswith('SI')]

# ═══════════════════════════════════════════════════════════════════════════
# Build presentation
# ═══════════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
_textbox(sl, Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.8),
         "Equivalent Roughness for Overland Flow\non Heterogeneous Hillslopes",
         sz=38, bold=True, align=PP_ALIGN.CENTER)
_textbox(sl, Inches(0.8), Inches(4.2), Inches(11.5), Inches(1),
         "Methods, Results & Figures",
         sz=22, align=PP_ALIGN.CENTER, color=(100, 100, 100))

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Simulation Setup
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
_textbox(sl, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "Simulation Setup (SWOF)", sz=28, bold=True)

tf_l = _textbox(sl, Inches(0.5), Inches(1.3), Inches(5.5), Inches(5.5),
                "Domain & Physics", sz=20, bold=True)
_bullet(tf_l, "1-D hillslope, full Saint-Venant equations", sz=17)
_bullet(tf_l, "Binary roughness: n\u1D65 (vegetation) and n\u2092 (bare)", sz=17)
_bullet(tf_l, "Hydrograph matching \u2192 n\u2091 for each simulation", sz=17)
_bullet(tf_l, "Smooth bed topography", sz=17)

tf_r = _textbox(sl, Inches(6.8), Inches(1.3), Inches(5.5), Inches(5.5),
                "Parameter Space", sz=20, bold=True)
_bullet(tf_r, "Vegetation fraction f\u1D65: 0.1 \u2013 0.9", sz=17)
_bullet(tf_r, "Patch lengthscale \u03C3 (L\u1D65): multiple levels", sz=17)
_bullet(tf_r, "Storm intensity p: 2 \u2013 8 mm/hr", sz=17)
_bullet(tf_r, "Storm duration t\u1D63: 30 \u2013 120 min", sz=17)
_bullet(tf_r, "Pattern anisotropy: isotropic + anisotropic", sz=17)
_bullet(tf_r, "Hundreds of unique simulations", sz=17)

# ═══════════════════════════════════════════════════════════════════════════
# Section divider: Main Figures
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
_textbox(sl, Inches(1), Inches(3), Inches(11), Inches(1.5),
         "Main Text Figures", sz=32, bold=True, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# Figure slides — image sized as large as possible
# ═══════════════════════════════════════════════════════════════════════════
def add_figure_slide(prs, fig_id, filename, caption):
    slide = prs.slides.add_slide(BLANK)
    img_path = FIG_DIR / filename
    if not img_path.exists():
        _textbox(slide, Inches(2), Inches(3), Inches(9), Inches(1),
                 f"{fig_id}: {filename} (file not found)", sz=18)
        return

    # Title
    _textbox(slide, Inches(0.3), Inches(0.12), Inches(12.5), Inches(0.5),
             fig_id, sz=24, bold=True)

    from PIL import Image
    with Image.open(img_path) as im:
        w_px, h_px = im.size
    aspect = w_px / h_px

    margin = Inches(0.25)
    top_off = Inches(0.6)
    cap_reserve = Inches(0.75)
    max_w = prs.slide_width - 2 * margin
    max_h = prs.slide_height - top_off - cap_reserve

    if aspect > (max_w / max_h):
        img_w = max_w; img_h = int(max_w / aspect)
    else:
        img_h = max_h; img_w = int(max_h * aspect)

    left = int((prs.slide_width - img_w) / 2)
    top = int(top_off)
    slide.shapes.add_picture(str(img_path), left, top, img_w, img_h)

    # Caption
    ct = top + img_h + Emu(30000)
    ch = prs.slide_height - ct - Emu(30000)
    if ch < Inches(0.3): ch = Inches(0.3)
    tb = slide.shapes.add_textbox(Inches(0.4), ct, prs.slide_width - Inches(0.8), ch)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = caption; p.font.size = Pt(11); p.font.italic = True

for fid, fn, cap in main_figs:
    add_figure_slide(prs, fid, fn, cap)

# ── SI section ─────────────────────────────────────────────────────────────
if si_figs:
    sl = prs.slides.add_slide(BLANK)
    _textbox(sl, Inches(1), Inches(3), Inches(11), Inches(1.5),
             "Supplementary Information", sz=32, bold=True, align=PP_ALIGN.CENTER)
    for fid, fn, cap in si_figs:
        add_figure_slide(prs, fid, fn, cap)

# ═══════════════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════════════
prs.save(str(OUT))
print(f"Created {OUT}")
print(f"  Main figures: {len(main_figs)}  |  SI figures: {len(si_figs)}")
print(f"  Total slides:   {len(prs.slides)}")
